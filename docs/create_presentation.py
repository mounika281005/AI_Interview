from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.7))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.3), Inches(5.9), Inches(5.7))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(6)

    return slide

# ============= CREATE SLIDES =============

# Slide 1: Title Slide
slide = add_title_slide(prs, "AI Mock Interview System", "An AI-Powered Mock Interview Platform with NLP-Based Evaluation")
# Add team info
info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(3))
tf = info_box.text_frame
tf.word_wrap = True
lines = [
    "Project Team:",
    "[Team Member 1] - [Reg. No. 1]        [Team Member 2] - [Reg. No. 2]",
    "[Team Member 3] - [Reg. No. 3]        [Team Member 4] - [Reg. No. 4]",
    "",
    "Project Guide: [Guide Name]",
    "Institute: [Institute/School Name]"
]
for i, line in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
        p.font.bold = True
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(51, 51, 51)

# Slide 2: Introduction
add_two_column_slide(prs, "Introduction",
    [
        "Project Overview",
        "",
        "AI Mock Interview System is a comprehensive",
        "AI-powered platform that helps candidates",
        "practice and improve their interview skills:",
        "",
        "• Voice-based interviews with real-time recording",
        "• NLP-powered evaluation for accurate assessment",
        "• Personalized feedback with improvement suggestions",
        "• Support for 40+ technologies and skills",
    ],
    [
        "Technologies Used",
        "",
        "Frontend: React.js, React Router, CSS3",
        "",
        "Backend: Python, FastAPI, SQLAlchemy",
        "",
        "AI/ML: OpenAI Whisper, Sentence Transformers,",
        "           NLTK, BERT",
        "",
        "Database: SQLite (Dev), PostgreSQL (Prod)",
        "",
        "Infrastructure: Docker, Uvicorn"
    ]
)

# Slide 3: Motivation
add_content_slide(prs, "Motivation",
    [
        "Why This Project?",
        "",
        "• Growing Demand: Competitive job market requires extensive interview practice",
        "• Cost Barrier: Traditional mock interviews are expensive and time-consuming",
        "• Limited Access: Shortage of professional interview coaches",
        "",
        "Gap in Existing Solutions:",
        "",
        "• Most platforms lack real-time voice-based interaction",
        "• Limited AI-powered feedback mechanisms",
        "• No comprehensive skill-based question generation",
        "",
        "Problems It Solves:",
        "",
        "• Provides 24/7 accessible interview practice",
        "• Delivers instant, objective feedback on responses",
        "• Offers skill-specific question generation for 40+ technologies"
    ]
)

# Slide 4: Objectives
add_content_slide(prs, "Objectives",
    [
        "Primary Objectives:",
        "",
        "1. Develop a voice-based mock interview platform that allows candidates",
        "   to practice interviews by speaking naturally, simulating real conditions",
        "",
        "2. Implement NLP-based answer evaluation using advanced AI models to",
        "   assess grammar, fluency, structure, and content relevance",
        "",
        "3. Create a comprehensive feedback system that provides actionable",
        "   improvement suggestions and ideal answer comparisons",
        "",
        "",
        "Validation Metrics:",
        "",
        "• Speech Recognition Accuracy: > 95% using Whisper",
        "• Answer Evaluation: Scoring correlation > 0.8 with expert ratings",
        "• User Experience: Complete interview flow in under 15 minutes"
    ]
)

# Slide 5: Background
add_two_column_slide(prs, "Background - Key Concepts",
    [
        "Speech-to-Text (STT)",
        "• Converting spoken audio into written text",
        "• Uses OpenAI Whisper model",
        "• Supports multiple languages and accents",
        "",
        "Natural Language Processing (NLP)",
        "• Analyzing and understanding human language",
        "• Evaluating grammar, fluency, coherence",
        "• Semantic similarity with ideal answers",
    ],
    [
        "Sentence Embeddings",
        "• Converting sentences into numerical vectors",
        "• Enables similarity comparison",
        "• Uses Sentence Transformers (BERT-based)",
        "",
        "Interview Assessment Criteria",
        "• Grammar: Correctness of sentence structure",
        "• Fluency: Natural flow and coherence",
        "• Structure: Logical organization (STAR method)",
        "• Relevance: Alignment with expected keywords"
    ]
)

# Slide 6: Literature Survey
add_content_slide(prs, "Literature Survey",
    [
        "Related Work:",
        "",
        "1. \"Automated Essay Scoring Using NLP\" - BERT-based Assessment Systems",
        "   • Demonstrated effectiveness of transformer models for text evaluation",
        "   • Limitation: Limited to written text, no speech component",
        "",
        "2. \"Robust Speech Recognition\" - OpenAI Whisper (2022)",
        "   • State-of-art speech recognition with high accuracy",
        "   • Limitation: Requires GPU for real-time processing",
        "",
        "3. \"Sentence-BERT\" - Reimers & Gurevych (2019)",
        "   • Efficient sentence embeddings for similarity matching",
        "   • Limitation: May miss contextual nuances",
        "",
        "Research Gap:",
        "No existing system combines voice-based interviews + NLP evaluation +",
        "skill-specific questions in a unified platform"
    ]
)

# Slide 7: Problem Identification
add_two_column_slide(prs, "Problem Identification & Proposed Solution",
    [
        "Problems Identified:",
        "",
        "1. Manual interview practice is expensive",
        "   → Limits practice opportunities",
        "",
        "2. Text-based platforms don't simulate",
        "   real interviews",
        "   → Poor preparation for actual interviews",
        "",
        "3. Generic feedback lacks actionable insights",
        "   → Slow improvement progress",
        "",
        "4. No skill-specific question generation",
        "   → Misaligned practice focus"
    ],
    [
        "Our Solution:",
        "",
        "✓ Free, unlimited AI-powered interviews",
        "",
        "✓ Voice recording with real-time",
        "   transcription",
        "",
        "✓ Per-question scoring with specific",
        "   improvement tips",
        "",
        "✓ Skill-specific questions for 40+",
        "   technologies",
        "",
        "✓ Progress tracking over multiple sessions"
    ]
)

# Slide 8: System Architecture
add_content_slide(prs, "System Architecture",
    [
        "                                    High-Level System Flow",
        "",
        "    ┌─────────────┐      ┌─────────────┐      ┌─────────────┐      ┌─────────────┐",
        "    │   SELECT    │ ──▶  │   ANSWER    │ ──▶  │     AI      │ ──▶  │    VIEW     │",
        "    │    SKILL    │      │  QUESTIONS  │      │ EVALUATION  │      │   RESULTS   │",
        "    └─────────────┘      └─────────────┘      └─────────────┘      └─────────────┘",
        "",
        "Architecture Layers:",
        "",
        "• Client Layer: React Web Application (UI/UX)",
        "",
        "• API Layer: FastAPI Backend (Interview Router, Question Router, Feedback Router)",
        "",
        "• AI Services Layer: Speech-to-Text (Whisper), NLP Evaluation (BERT), Scoring Engine",
        "",
        "• Data Layer: SQLite/PostgreSQL + File Storage"
    ]
)

# Slide 9: Dataset Description
add_two_column_slide(prs, "Dataset Description",
    [
        "Technologies Covered (40+ Skills):",
        "",
        "Programming Languages:",
        "Python, JavaScript, Java, C++, Go, Rust",
        "",
        "Frontend:",
        "React, Angular, Vue.js, HTML/CSS",
        "",
        "Backend:",
        "Node.js, Django, FastAPI, Spring Boot",
        "",
        "Databases:",
        "SQL, MongoDB, PostgreSQL, Redis",
        "",
        "Cloud & DevOps:",
        "AWS, Azure, GCP, Docker, Kubernetes"
    ],
    [
        "Question Database:",
        "",
        "• Total Questions: 500+ unique questions",
        "• Question Types: Technical, Behavioral,",
        "  Situational",
        "• Difficulty: Beginner, Intermediate, Advanced",
        "• Per Interview: 5 dynamically selected",
        "",
        "Data Generated Per Interview:",
        "",
        "• Audio Files: WebM format recordings",
        "• Transcriptions: Text from STT",
        "• Scores: 4 criteria scores (0-5 each)",
        "• Feedback: AI-generated improvement tips"
    ]
)

# Slide 10: AI Approach
add_content_slide(prs, "AI Approach & Methodology",
    [
        "AI Processing Pipeline:",
        "",
        "Stage 1: Speech-to-Text → OpenAI Whisper converts audio to text",
        "",
        "Stage 2: Text Analysis → NLTK & SpaCy analyze grammar and fluency",
        "",
        "Stage 3: Semantic Evaluation → Sentence-BERT computes similarity scores",
        "",
        "Stage 4: Feedback Generation → Personalized improvement suggestions",
        "",
        "",
        "Scoring Criteria (Each 25% weight):",
        "",
        "• Grammar: Grammatical correctness (NLTK grammar checker)",
        "• Fluency: Natural language flow (Perplexity analysis)",
        "• Structure: Logical organization (STAR pattern matching)",
        "• Similarity: Relevance to ideal answer (Cosine similarity)"
    ]
)

# Slide 11: Preliminary Results
add_two_column_slide(prs, "Preliminary Results",
    [
        "System Performance Metrics:",
        "",
        "• Speech Recognition Accuracy: [To be updated]",
        "  Target: > 95%",
        "",
        "• Average Response Time: [To be updated]",
        "  Target: < 3 seconds",
        "",
        "• Scoring Consistency: [To be updated]",
        "  Target: > 0.8 correlation",
        "",
        "",
        "Grade Scale:",
        "A+: 90-100%    B+: 70-79%    C: 50-59%",
        "A: 80-89%      B: 60-69%     D: 40-49%"
    ],
    [
        "Sample Evaluation Output:",
        "",
        "Question: \"Explain closures in JavaScript\"",
        "",
        "Scores:",
        "  • Grammar: 4.2/5",
        "  • Fluency: 3.8/5",
        "  • Structure: 4.0/5",
        "  • Similarity: 4.5/5",
        "",
        "Total: 82.5% (Grade: A)",
        "",
        "Feedback:",
        "  ✓ Clear explanation of core concept",
        "  ✗ Include practical code example"
    ]
)

# Slide 12: Demo
add_content_slide(prs, "Demo",
    [
        "Live Demonstration Flow:",
        "",
        "1. Start Interview",
        "   • Select technology (e.g., \"JavaScript\")",
        "   • System generates 5 interview questions",
        "",
        "2. Answer Questions",
        "   • Question displayed on screen",
        "   • User clicks \"Record\" and speaks answer",
        "   • Audio recorded via browser microphone",
        "",
        "3. Submit for Evaluation",
        "   • All 5 recordings submitted to AI pipeline",
        "   • Each answer processed through STT → NLP → Scoring",
        "",
        "4. View Results",
        "   • Per-question scores and feedback displayed",
        "   • Strengths and improvements highlighted",
        "   • Overall grade and progress tracking"
    ]
)

# Slide 13: Team Contribution
add_content_slide(prs, "Team Contribution",
    [
        "",
        "┌──────────────────────┬────────────────────────────────────────────────┐",
        "│     Team Member      │                  Contributions                 │",
        "├──────────────────────┼────────────────────────────────────────────────┤",
        "│  [Member 1 Name]     │  Frontend Development (React), UI/UX Design   │",
        "├──────────────────────┼────────────────────────────────────────────────┤",
        "│  [Member 2 Name]     │  Backend Development (FastAPI), API Design    │",
        "├──────────────────────┼────────────────────────────────────────────────┤",
        "│  [Member 3 Name]     │  AI/ML Pipeline (STT, NLP Evaluation)         │",
        "├──────────────────────┼────────────────────────────────────────────────┤",
        "│  [Member 4 Name]     │  Database Design, Testing, Documentation      │",
        "└──────────────────────┴────────────────────────────────────────────────┘",
        "",
        "",
        "Note: Update team member names and specific contributions"
    ]
)

# Slide 14: Tools & Technologies
add_two_column_slide(prs, "Tools & Technologies",
    [
        "Development Stack:",
        "",
        "Frontend:",
        "• React.js - User interface",
        "• React Router - Client-side routing",
        "• CSS3 - Styling",
        "",
        "Backend:",
        "• Python 3.10+ - Server programming",
        "• FastAPI - REST API framework",
        "• SQLAlchemy - ORM",
        "• Pydantic - Data validation",
        "",
        "Database:",
        "• SQLite - Development",
        "• PostgreSQL - Production"
    ],
    [
        "AI/ML Technologies:",
        "",
        "• OpenAI Whisper - Speech-to-Text",
        "• Sentence Transformers - Embeddings",
        "• NLTK - Grammar analysis",
        "• SpaCy - NLP processing",
        "",
        "DevOps & Tools:",
        "",
        "• Docker - Containerization",
        "• Uvicorn - ASGI server",
        "• Git/GitHub - Version control",
        "• VS Code - Development IDE",
        "• Postman - API testing"
    ]
)

# Slide 15: References
add_content_slide(prs, "References",
    [
        "[1] A. Radford et al., \"Robust Speech Recognition via Large-Scale Weak",
        "    Supervision,\" OpenAI, 2022. Available: https://openai.com/research/whisper",
        "",
        "[2] N. Reimers and I. Gurevych, \"Sentence-BERT: Sentence Embeddings using",
        "    Siamese BERT-Networks,\" EMNLP, 2019.",
        "",
        "[3] S. Bird, E. Klein, and E. Loper, Natural Language Processing with Python,",
        "    O'Reilly Media, 2009.",
        "",
        "[4] S. Ramírez, \"FastAPI Documentation,\" 2023. Available: https://fastapi.tiangolo.com/",
        "",
        "[5] React Team, \"React Documentation,\" Meta, 2023. Available: https://react.dev/",
        "",
        "[6] J. Devlin et al., \"BERT: Pre-training of Deep Bidirectional Transformers,\"",
        "    NAACL-HLT, 2019.",
        "",
        "[7] M. Honnibal and I. Montani, \"spaCy 2: Natural Language Understanding,\" 2017."
    ]
)

# Slide 16: Thank You
slide = add_title_slide(prs, "Thank You!", "Questions?")
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(2))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "Contact: [team-email@example.com]"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(51, 51, 51)
p = tf.add_paragraph()
p.text = "GitHub: [github.com/your-repo]"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(51, 51, 51)

# Save the presentation
prs.save('f:/Interview/AI_Interview/docs/AI_Mock_Interview_Presentation.pptx')
print("Presentation created successfully!")
print("Location: f:/Interview/AI_Interview/docs/AI_Mock_Interview_Presentation.pptx")
